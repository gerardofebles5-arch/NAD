"""
Fase 5: Multi-formato Input — Parseo de PDF, DOCX, PPTX, XLSX, HTML.
================================================================================
Convierte documentos en múltiples formatos a imágenes procesables
y extrae contenido estructurado.

Formatos soportados:
  - PDF (PyMuPDF)
  - DOCX (python-docx)
  - PPTX (python-pptx)
  - XLSX (openpyxl)
  - HTML (BeautifulSoup)
  - Imágenes (OpenCV)
"""

import os
import io
import re
from dataclasses import dataclass, field
from typing import List, Optional, Dict, Any
from pathlib import Path

import cv2
import numpy as np


@dataclass
class PageContent:
    """Contenido de una página del documento."""
    page_number: int
    image: Optional[np.ndarray] = None   # Imagen BGR
    text: str = ""                       # Texto extraído
    tables: List[Dict] = field(default_factory=list)  # Tablas extraídas
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class DocumentContent:
    """Contenido completo de un documento."""
    filename: str
    format: str
    pages: List[PageContent]
    metadata: Dict[str, Any] = field(default_factory=dict)

    @property
    def num_pages(self) -> int:
        return len(self.pages)

    def get_images(self) -> List[np.ndarray]:
        """Retorna solo las imágenes de las páginas."""
        return [p.image for p in self.pages if p.image is not None]

    def get_full_text(self) -> str:
        """Retorna todo el texto del documento."""
        return "\n\n".join(
            f"--- Página {p.page_number} ---\n{p.text}"
            for p in self.pages
            if p.text
        )


class DocumentParser:
    """
    Parsea documentos en múltiples formatos.

    Uso:
        parser = DocumentParser()
        doc = parser.parse("documento.pdf")
        for page in doc.pages:
            print(f"Página {page.page_number}: {len(page.text)} chars")
    """

    # Formatos soportados
    SUPPORTED_FORMATS = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'docx',
        '.pptx': 'pptx',
        '.ppt': 'pptx',
        '.xlsx': 'xlsx',
        '.xls': 'xlsx',
        '.html': 'html',
        '.htm': 'html',
        '.png': 'image',
        '.jpg': 'image',
        '.jpeg': 'image',
        '.bmp': 'image',
        '.tiff': 'image',
        '.tif': 'image',
    }

    def __init__(self, dpi: int = 150):
        """
        Args:
            dpi: Resolución para rendering de PDF
        """
        self.dpi = dpi

    def parse(self, file_path: str) -> DocumentContent:
        """
        Parsea un documento.

        Args:
            file_path: Ruta al archivo

        Returns:
            DocumentContent con todas las páginas
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        if ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Formato no soportado: {ext}")

        format_type = self.SUPPORTED_FORMATS[ext]

        print(f"  [Parser] Parseando {format_type}: {path.name}")

        if format_type == 'pdf':
            return self._parse_pdf(path)
        elif format_type == 'docx':
            return self._parse_docx(path)
        elif format_type == 'pptx':
            return self._parse_pptx(path)
        elif format_type == 'xlsx':
            return self._parse_xlsx(path)
        elif format_type == 'html':
            return self._parse_html(path)
        elif format_type == 'image':
            return self._parse_image(path)
        else:
            raise ValueError(f"Formato no implementado: {format_type}")

    def _parse_pdf(self, path: Path) -> DocumentContent:
        """Parsea PDF usando PyMuPDF."""
        pages = []

        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(path))
            print(f"  [Parser] PDF: {len(doc)} páginas")

            for i, page in enumerate(doc):
                # Renderizar página a imagen
                mat = fitz.Matrix(self.dpi / 72, self.dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                    pix.height, pix.width, pix.n
                )

                # Convertir a BGR si es necesario
                if pix.n == 4:  # RGBA
                    img = cv2.cvtColor(img, cv2.COLOR_RGBA2BGR)
                elif pix.n == 1:  # Grayscale
                    img = cv2.cvtColor(img, cv2.COLOR_GRAY2BGR)

                # Extraer texto
                text = page.get_text("text")

                pages.append(PageContent(
                    page_number=i + 1,
                    image=img,
                    text=text,
                    metadata={"width": pix.width, "height": pix.height},
                ))

            doc.close()

        except ImportError:
            print("  [Parser] PyMuPDF no instalado. Instale: pip install PyMuPDF")
            pages = self._fallback_pdf(path)

        return DocumentContent(
            filename=path.name,
            format="pdf",
            pages=pages,
            metadata={"dpi": self.dpi},
        )

    def _fallback_pdf(self, path: Path) -> List[PageContent]:
        """Fallback para PDF sin PyMuPDF."""
        try:
            import pdfplumber

            pages = []
            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    tables = page.extract_tables() or []

                    pages.append(PageContent(
                        page_number=i + 1,
                        text=text,
                        tables=[{"rows": t} for t in tables],
                    ))

            return pages

        except ImportError:
            print("  [Parser] pdfplumber no disponible")
            return []

    def _parse_docx(self, path: Path) -> DocumentContent:
        """Parsea DOCX usando python-docx."""
        pages = []

        try:
            from docx import Document

            doc = Document(str(path))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extraer tablas
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append({"rows": table_data})

            full_text = "\n".join(text_parts)

            # Convertir a imagen (placeholder - sin rendering real)
            img = self._text_to_image(full_text)

            pages.append(PageContent(
                page_number=1,
                image=img,
                text=full_text,
                tables=tables,
            ))

        except ImportError:
            print("  [Parser] python-docx no instalado. Instale: pip install python-docx")

        return DocumentContent(
            filename=path.name,
            format="docx",
            pages=pages,
        )

    def _parse_pptx(self, path: Path) -> DocumentContent:
        """Parsea PPTX usando python-pptx."""
        pages = []

        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation(str(path))

            for i, slide in enumerate(prs.slides):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)

                full_text = "\n".join(text_parts)
                img = self._text_to_image(full_text)

                pages.append(PageContent(
                    page_number=i + 1,
                    image=img,
                    text=full_text,
                ))

        except ImportError:
            print("  [Parser] python-pptx no instalado. Instale: pip install python-pptx")

        return DocumentContent(
            filename=path.name,
            format="pptx",
            pages=pages,
        )

    def _parse_xlsx(self, path: Path) -> DocumentContent:
        """Parsea XLSX usando openpyxl."""
        pages = []

        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(path), data_only=True)

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []

                for row in ws.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    rows.append(row_data)

                # Convertir a imagen
                text = "\n".join([" | ".join(row) for row in rows[:50]])
                img = self._text_to_image(text)

                pages.append(PageContent(
                    page_number=len(pages) + 1,
                    image=img,
                    text=text,
                    tables=[{"rows": rows}],
                    metadata={"sheet": sheet_name},
                ))

            wb.close()

        except ImportError:
            print("  [Parser] openpyxl no instalado. Instale: pip install openpyxl")

        return DocumentContent(
            filename=path.name,
            format="xlsx",
            pages=pages,
        )

    def _parse_html(self, path: Path) -> DocumentContent:
        """Parsea HTML usando BeautifulSoup."""
        pages = []

        try:
            from bs4 import BeautifulSoup

            with open(path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, 'html.parser')

            # Extraer texto
            text = soup.get_text(separator='\n', strip=True)

            # Extraer tablas
            tables = []
            for table in soup.find_all('table'):
                rows = []
                for tr in table.find_all('tr'):
                    row = [td.get_text(strip=True) for td in tr.find_all(['td', 'th'])]
                    rows.append(row)
                if rows:
                    tables.append({"rows": rows})

            img = self._text_to_image(text)

            pages.append(PageContent(
                page_number=1,
                image=img,
                text=text,
                tables=tables,
            ))

        except ImportError:
            print("  [Parser] BeautifulSoup no instalado. Instale: pip install beautifulsoup4")

        return DocumentContent(
            filename=path.name,
            format="html",
            pages=pages,
        )

    def _parse_image(self, path: Path) -> DocumentContent:
        """Parsea imagen directamente."""
        img = cv2.imread(str(path))

        if img is None:
            raise ValueError(f"No se pudo leer la imagen: {path}")

        return DocumentContent(
            filename=path.name,
            format="image",
            pages=[PageContent(
                page_number=1,
                image=img,
                metadata={"width": img.shape[1], "height": img.shape[0]},
            )],
        )

    def _text_to_image(self, text: str, max_width: int = 800) -> np.ndarray:
        """Convierte texto a imagen básica (placeholder para documentos sin imagen)."""
        lines = text.split('\n')[:50]  # Máximo 50 líneas
        line_height = 20
        height = max(100, len(lines) * line_height + 40)

        img = np.ones((height, max_width, 3), dtype=np.uint8) * 255

        for i, line in enumerate(lines[:height // line_height]):
            y = 30 + i * line_height
            # Truncar línea larga
            display_line = line[:60]
            cv2.putText(
                img, display_line, (20, y),
                cv2.FONT_HERSHEY_SIMPLEX, 0.4, (0, 0, 0), 1
            )

        return img


# ═══════════════════════════════════════════════════════════════
#  Función de conveniencia
# ═══════════════════════════════════════════════════════════════

def parse_document(file_path: str, dpi: int = 150) -> DocumentContent:
    """Función de conveniencia para parsear documentos."""
    parser = DocumentParser(dpi=dpi)
    return parser.parse(file_path)
